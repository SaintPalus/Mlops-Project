"""Generate Food Waste Prediction PowerPoint — concise 10-minute version.
   10 slides · Output: Food_Waste_Prediction_Slides.pptx
"""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

# ── Colour Palette ─────────────────────────────────────────────────────────────
C_BG     = RGBColor(0x0D, 0x1B, 0x2A)
C_ACCENT = RGBColor(0x00, 0xC9, 0xA7)
C_YELLOW = RGBColor(0xFF, 0xD1, 0x66)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT  = RGBColor(0xB0, 0xC4, 0xDE)
C_CARD   = RGBColor(0x16, 0x2D, 0x44)
C_GREEN  = RGBColor(0x4C, 0xAF, 0x50)
C_ORANGE = RGBColor(0xFF, 0x98, 0x00)
C_RED    = RGBColor(0xEF, 0x53, 0x50)

# ── Setup ──────────────────────────────────────────────────────────────────────
prs  = Presentation()
prs.slide_width  = Emu(int(13.33 * 914400))
prs.slide_height = Emu(int(7.5  * 914400))
BLANK = prs.slide_layouts[6]

def add_bg(slide):
    s = slide.shapes.add_shape(1,0,0, prs.slide_width, prs.slide_height)
    s.fill.solid(); s.fill.fore_color.rgb = C_BG
    s.line.fill.background()

def txb(slide, text, x, y, w, h, size=20, bold=False, color=C_WHITE,
        align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tb.word_wrap = True
    tf = tb.text_frame; tf.word_wrap = True
    p  = tf.paragraphs[0]; p.alignment = align
    r  = p.add_run(); r.text = text
    r.font.size=Pt(size); r.font.bold=bold; r.font.italic=italic
    r.font.color.rgb=color; r.font.name="Sarabun"
    return tb

def box(slide, x, y, w, h, fill=C_CARD, line=None, lw=Pt(1), radius=False):
    t = 5 if radius else 1
    s = slide.shapes.add_shape(t, Inches(x), Inches(y), Inches(w), Inches(h))
    s.fill.solid(); s.fill.fore_color.rgb = fill
    if line: s.line.color.rgb=line; s.line.width=lw
    else: s.line.fill.background()
    return s

def bar(slide, title, subtitle=None):
    """Accent bar + title line."""
    box(slide, 0.4, 1.1, 1.2, 0.07, fill=C_ACCENT)
    txb(slide, title, 0.4, 0.2, 12.5, 0.75, size=30, bold=True)
    if subtitle:
        txb(slide, subtitle, 0.4, 0.88, 12.5, 0.4, size=14, italic=True, color=C_LIGHT)

def div(slide, y):
    box(slide, 0.4, y, 12.5, 0.02, fill=C_ACCENT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)

box(s, 9.6, 0.2, 3.3, 3.3, fill=C_CARD, radius=True)
txb(s, "🍱", 10.2, 0.5, 2.2, 2.0, size=80, align=PP_ALIGN.CENTER)

box(s, 0.4, 0.4, 2.8, 0.55, fill=C_ACCENT, radius=True)
txb(s, "SDG 2  ·  SDG 12", 0.45, 0.43, 2.7, 0.5,
    size=15, bold=True, color=C_BG, align=PP_ALIGN.CENTER)

txb(s, "ระบบพยากรณ์ความต้องการอาหาร", 0.4, 1.15, 9.0, 0.95, size=40, bold=True)
txb(s, "และวางแผนลดขยะ", 0.4, 2.0, 9.0, 0.85, size=40, bold=True, color=C_ACCENT)
txb(s, "Food Waste Prediction & Optimization Planner", 0.4, 2.88, 9.0, 0.55,
    size=19, italic=True, color=C_LIGHT)

div(s, 3.65)
txb(s, "Palus Kaewaram  (66070131)  ·  Natthapol Aiemburanont  (66070062)",
    0.4, 3.82, 12.0, 0.45, size=16, color=C_LIGHT)
txb(s, "วิชา MLOps  |  มีนาคม 2569", 0.4, 4.28, 8.0, 0.4, size=15, color=C_LIGHT)

for i, t in enumerate(["PyTorch","Prefect","MLflow","FastAPI","Streamlit","Docker"]):
    px = 0.4 + i * 2.1
    box(s, px, 6.4, 1.95, 0.55, fill=C_CARD, radius=True, line=C_ACCENT, lw=Pt(1))
    txb(s, t, px, 6.4, 1.95, 0.55, size=15, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
bar(s, "ปัญหา  —  Food Waste", "1 ใน 3 ของอาหารทั่วโลกถูกทิ้งเป็นขยะ  (FAO)")

stats = [
    (C_ACCENT,  "1 ใน 3",       "ของอาหารโลกถูกทิ้ง\nทุกปี  (FAO)"),
    (C_RED,     "733 ล้านคน",   "ยังขาดแคลนอาหาร\n(SOFI 2023)"),
    (C_ORANGE,  "1.3 พันล้านตัน","สูญเสียโดยเปล่าประโยชน์\nต่อปี"),
]
for i,(col,val,desc) in enumerate(stats):
    px = 0.4 + i * 4.3
    box(s, px, 1.35, 4.1, 2.3, fill=C_CARD, radius=True, line=col, lw=Pt(2))
    txb(s, val,  px+0.1, 1.45, 3.9, 0.9, size=36, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, desc, px+0.1, 2.3,  3.9, 1.1, size=16, color=C_LIGHT, align=PP_ALIGN.CENTER)

div(s, 3.85)
txb(s, "ปัญหาที่ร้านเบเกอรี่ / โรงอาหารเจอ", 0.4, 4.0, 8.0, 0.45, size=18, bold=True, color=C_YELLOW)

probs = [
    (C_YELLOW, "เตรียมเกิน",    "ไม่รู้ว่าแต่ละวันต้องเตรียมเท่าไร"),
    (C_RED,    "ตัดสินใจด้วยความรู้สึก", "ไม่มีข้อมูลสนับสนุน"),
    (C_ORANGE, "ยอดขายผันผวน",  "วันหยุด ≠ วันธรรมดา"),
    (C_GREEN,  "ขาดเครื่องมือ", "ซอฟต์แวร์ Enterprise แพงเกินไป"),
]
for i,(col,t,d) in enumerate(probs):
    px = 0.4 + i * 3.24
    box(s, px, 4.55, 3.08, 2.1, fill=C_CARD, radius=True, line=col, lw=Pt(1))
    txb(s, t, px+0.15, 4.65, 2.78, 0.45, size=15, bold=True, color=col)
    txb(s, d, px+0.15, 5.15, 2.78, 1.2,  size=13, color=C_LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
bar(s, "Solution  —  แนวทางแก้ไข")

steps = [
    (C_ACCENT,  "1", "โหลดข้อมูล",   "20,507 transactions จริง"),
    (C_YELLOW,  "2", "Feature Eng.", "4 features: วัน/เดือน/วันที่"),
    (C_GREEN,   "3", "MLP Training", "4→32→16→1  ·  1,000 epochs"),
    (C_ORANGE,  "4", "Safety Stock", "เพิ่ม +5% อัตโนมัติ"),
    (C_RED,     "5", "Dashboard",   "Streamlit ภาษาไทย"),
]
for i,(col,num,title,desc) in enumerate(steps):
    px = 0.4 + i * 2.57
    box(s, px+0.55, 1.35, 0.85, 0.85, fill=col, radius=True)
    txb(s, num, px+0.55, 1.35, 0.85, 0.85, size=24, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    if i < 4:
        txb(s, "→", px+1.55, 1.5, 0.55, 0.55, size=22, color=C_LIGHT, align=PP_ALIGN.CENTER)
    box(s, px, 2.35, 2.45, 2.2, fill=C_CARD, radius=True)
    txb(s, title, px+0.1, 2.47, 2.25, 0.5, size=16, bold=True, color=col)
    txb(s, desc,  px+0.1, 3.0,  2.25, 1.3, size=13, color=C_LIGHT)

div(s, 4.75)
txb(s, "นวัตกรรมหลัก  :  Per-Item Mean Normalization", 0.4, 4.9, 9.0, 0.5,
    size=18, bold=True, color=C_YELLOW)

inno = [
    (C_ACCENT, "ปัญหา",   "Coffee ~35/วัน  vs  Brownie ~5/วัน → scale ต่างกัน 7x"),
    (C_GREEN,  "วิธีแก้", "demand_norm = demand / item_mean  →  ทุก item normalize รอบ 1.0"),
    (C_ORANGE, "Inference","predicted = norm_factor × item_mean  →  แปลงกลับสู่หน่วยจริง"),
]
for i,(col,t,d) in enumerate(inno):
    py = 5.5 + i * 0.56
    box(s, 0.4, py, 12.5, 0.5, fill=C_CARD, radius=True)
    txb(s, t, 0.55, py+0.07, 1.4, 0.36, size=13, bold=True, color=col)
    txb(s, d, 2.0,  py+0.07, 10.7, 0.36, size=13, color=C_LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — System Architecture
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
bar(s, "System Architecture  —  สถาปัตยกรรมระบบ")

# Training pipeline (left box)
box(s, 0.35, 1.35, 5.8, 5.75, fill=RGBColor(0x10,0x22,0x34), radius=True, line=C_ACCENT, lw=Pt(1))
txb(s, "Training Pipeline", 0.5, 1.42, 5.5, 0.45, size=16, bold=True, color=C_ACCENT)

pipe = [
    (C_ACCENT,  "📂  Data",           "bakery_sales_revised.csv  (20,507 rows)"),
    (C_YELLOW,  "⚙️  Prefect",        "Orchestrate @flow / @task"),
    (C_GREEN,   "🧠  MLP Training",   "4→32→16→1  ·  Adam  ·  L1Loss"),
    (C_ORANGE,  "📊  MLflow+DagsHub", "Log params & MAE"),
    (C_LIGHT,   "💾  Save Artifacts", "demand_model.pth  ·  item_map.json"),
]
for i,(col,t,sub) in enumerate(pipe):
    py = 2.0 + i * 0.97
    box(s, 0.55, py, 5.4, 0.8, fill=C_CARD, radius=True)
    txb(s, t,   0.72, py+0.07, 2.3, 0.36, size=14, bold=True, color=col)
    txb(s, sub, 3.05, py+0.07, 2.75, 0.6, size=12, color=C_LIGHT)
    if i<4:
        txb(s, "↓", 2.8, py+0.82, 0.5, 0.15, size=13, color=C_ACCENT, align=PP_ALIGN.CENTER)

# Arrow
txb(s, "→", 6.2, 3.85, 0.7, 0.7, size=32, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
txb(s, "model\nartifacts", 6.05, 4.5, 1.0, 0.6, size=11, color=C_LIGHT, align=PP_ALIGN.CENTER)

# Serving (right box)
box(s, 7.1, 1.35, 6.0, 5.75, fill=RGBColor(0x10,0x22,0x34), radius=True, line=C_YELLOW, lw=Pt(1))
txb(s, "Serving Layer", 7.25, 1.42, 5.7, 0.45, size=16, bold=True, color=C_YELLOW)

box(s, 7.3, 2.05, 2.6, 2.0, fill=C_CARD, radius=True, line=C_GREEN, lw=Pt(1))
txb(s, "FastAPI Backend", 7.45, 2.12, 2.3, 0.42, size=14, bold=True, color=C_GREEN)
txb(s, "Render.com", 7.45, 2.52, 2.3, 0.32, size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)
for j,ep in enumerate(["/health","/items","/predict"]):
    txb(s, ep, 7.45, 2.88+j*0.35, 2.3, 0.32, size=12, color=C_LIGHT)

txb(s, "↔", 10.0, 2.85, 0.6, 0.6, size=20, color=C_YELLOW, align=PP_ALIGN.CENTER)

box(s, 10.7, 2.05, 2.15, 2.0, fill=C_CARD, radius=True, line=C_ORANGE, lw=Pt(1))
txb(s, "Streamlit\nFrontend", 10.75, 2.12, 2.05, 0.7, size=13, bold=True, color=C_ORANGE, align=PP_ALIGN.CENTER)
txb(s, "Streamlit Cloud\nภาษาไทย", 10.75, 2.85, 2.05, 0.85, size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)

box(s, 7.3, 4.22, 5.6, 1.2, fill=C_CARD, radius=True, line=C_ACCENT, lw=Pt(1))
txb(s, "MLflow  +  DagsHub", 7.45, 4.3, 5.3, 0.4, size=14, bold=True, color=C_ACCENT)
txb(s, "Experiment tracking  ·  Model registry  ·  Run comparison", 7.45, 4.72, 5.3, 0.55, size=12, color=C_LIGHT)

box(s, 7.3, 5.6, 5.6, 1.15, fill=C_CARD, radius=True)
txb(s, "Docker Compose  —  3 Services",
    7.45, 5.67, 5.3, 0.4, size=13, bold=True, color=C_LIGHT)
txb(s, "mlflow :5000  ·  backend :8000  ·  frontend :8501",
    7.45, 6.06, 5.3, 0.38, size=12, color=C_LIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — ML Model
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
bar(s, "ML Model  —  MLP Neural Network",
    "4 Features  ·  4→32→16→1  ·  Adam  ·  L1Loss  ·  1,000 epochs")

# Architecture visual (left)
box(s, 0.35, 1.35, 6.2, 5.75, fill=RGBColor(0x10,0x22,0x34), radius=True, line=C_ACCENT, lw=Pt(1))
txb(s, "Architecture", 0.5, 1.42, 5.9, 0.45, size=16, bold=True, color=C_ACCENT)

# Input
box(s, 0.55, 2.05, 1.65, 0.55, fill=C_CARD, radius=True, line=C_ACCENT)
txb(s, "Input (4)", 0.55, 2.05, 1.65, 0.55, size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)
feats = [("day_of_week",C_ACCENT),("is_weekend",C_YELLOW),("month_norm",C_GREEN),("day_norm",C_ORANGE)]
for i,(f,fc) in enumerate(feats):
    py = 2.72 + i*0.5
    box(s, 0.55, py, 1.65, 0.42, fill=C_CARD, radius=True, line=fc)
    txb(s, f, 0.55, py, 1.65, 0.42, size=12, color=fc, align=PP_ALIGN.CENTER)

txb(s, "→", 2.3, 3.4, 0.5, 0.5, size=22, color=C_LIGHT, align=PP_ALIGN.CENTER)
box(s, 2.9, 3.0, 1.2, 1.0, fill=C_CARD, radius=True, line=C_YELLOW)
txb(s, "32\nReLU", 2.9, 3.0, 1.2, 1.0, size=20, bold=True, color=C_YELLOW, align=PP_ALIGN.CENTER)
txb(s, "→", 4.2, 3.38, 0.5, 0.5, size=22, color=C_LIGHT, align=PP_ALIGN.CENTER)
box(s, 4.8, 3.1, 1.05, 0.8, fill=C_CARD, radius=True, line=C_GREEN)
txb(s, "16\nReLU", 4.8, 3.1, 1.05, 0.8, size=18, bold=True, color=C_GREEN, align=PP_ALIGN.CENTER)
txb(s, "→", 5.95, 3.38, 0.5, 0.5, size=22, color=C_LIGHT, align=PP_ALIGN.CENTER)
box(s, 3.5, 4.25, 1.7, 0.55, fill=C_CARD, radius=True, line=C_ORANGE)
txb(s, "Output (1)\ndemand_norm", 3.4, 4.28, 2.1, 0.5, size=12, bold=True, color=C_ORANGE)

box(s, 0.55, 5.25, 5.8, 1.5, fill=C_CARD, radius=True)
txb(s, "Normalization:", 0.7, 5.32, 5.5, 0.38, size=13, bold=True, color=C_YELLOW)
txb(s, "demand_norm  =  demand / item_mean", 0.7, 5.7, 5.5, 0.35, size=13, color=C_WHITE)
txb(s, "predicted    =  norm_factor × item_mean", 0.7, 6.05, 5.5, 0.35, size=13, color=C_WHITE)

# Training config (right)
box(s, 6.8, 1.35, 6.1, 5.75, fill=RGBColor(0x10,0x22,0x34), radius=True, line=C_YELLOW, lw=Pt(1))
txb(s, "Training Config", 6.95, 1.42, 5.8, 0.45, size=16, bold=True, color=C_YELLOW)

cfg = [("Optimizer","Adam"),("Learning Rate","0.005"),("Epochs","1,000"),
       ("Loss","L1Loss (MAE)"),("Batch","Full batch"),("Device","CPU")]
for i,(k,v) in enumerate(cfg):
    py = 2.05 + i*0.42
    txb(s, k+" :", 6.95, py, 2.5, 0.38, size=14, color=C_LIGHT)
    txb(s, v, 9.5, py, 3.2, 0.38, size=14, bold=True)

txb(s, "Per-Item Mean Demand", 6.95, 4.65, 5.8, 0.42, size=15, bold=True, color=C_ACCENT)
items2 = [("Coffee","34.6"),("Tea","9.1"),("Bread","20.9"),("Sandwich","5.0"),
          ("Cake","7.0"),("Pastry","5.7"),("Medialuna","4.7"),("Cookies","3.8"),
          ("Hot Choc.","4.0"),("Brownie","4.5")]
for i,(it,m) in enumerate(items2):
    ci=i%2; ri=i//2
    px=6.95+ci*3.0; py=5.15+ri*0.38
    txb(s, it,  px,     py, 2.2, 0.34, size=12, color=C_LIGHT)
    txb(s, m,   px+2.2, py, 0.75, 0.34, size=12, bold=True, color=C_YELLOW, align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MLOps Pipeline
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
bar(s, "MLOps Pipeline  —  Prefect + MLflow + DagsHub")

# Prefect (left)
box(s, 0.35, 1.35, 6.0, 5.75, fill=RGBColor(0x10,0x22,0x34), radius=True, line=C_ACCENT, lw=Pt(1))
txb(s, "Prefect Orchestration", 0.5, 1.42, 5.7, 0.45, size=16, bold=True, color=C_ACCENT)

tasks = [
    (C_ACCENT,  "@flow",            "food-waste-training-pipeline"),
    (C_YELLOW,  "@task  load_data", "อ่าน CSV · 20,507 rows"),
    (C_GREEN,   "@task  train",     "PyTorch MLP · Adam · L1Loss"),
    (C_ORANGE,  "@task  mlflow_log","log params + MAE_normalized"),
    (C_LIGHT,   "@task  save",      "demand_model.pth · item_map.json"),
]
for i,(col,t,d) in enumerate(tasks):
    py = 2.05 + i*0.93
    box(s, 0.55, py, 5.6, 0.78, fill=C_CARD, radius=True)
    txb(s, t, 0.72, py+0.05, 2.5, 0.36, size=14, bold=True, color=col)
    txb(s, d, 3.25, py+0.05, 2.7, 0.6, size=12, color=C_LIGHT)
    if i<4:
        txb(s, "↓", 3.0, py+0.8, 0.5, 0.15, size=12, color=C_ACCENT, align=PP_ALIGN.CENTER)

# MLflow (right)
box(s, 6.55, 1.35, 6.4, 5.75, fill=RGBColor(0x10,0x22,0x34), radius=True, line=C_YELLOW, lw=Pt(1))
txb(s, "MLflow  +  DagsHub", 6.7, 1.42, 6.1, 0.45, size=16, bold=True, color=C_YELLOW)

box(s, 6.75, 2.05, 6.0, 3.0, fill=C_CARD, radius=True)
txb(s, "Parameters Logged", 6.9, 2.12, 5.7, 0.4, size=14, bold=True, color=C_YELLOW)
params = [("epochs","1000"),("lr","0.005"),("architecture","4-32-16-1"),
          ("normalization","per_item_mean"),("loss","L1Loss")]
for i,(k,v) in enumerate(params):
    py = 2.6 + i*0.38
    txb(s, k+" :", 6.9, py, 2.5, 0.34, size=13, color=C_LIGHT)
    txb(s, v,     9.45, py, 3.0, 0.34, size=13, bold=True)

box(s, 6.75, 5.2, 6.0, 1.0, fill=C_CARD, radius=True)
txb(s, "Metric", 6.9, 5.27, 2.0, 0.36, size=13, bold=True, color=C_GREEN)
txb(s, "MAE_normalized  ~0.10 – 0.12", 6.9, 5.62, 5.7, 0.35, size=13, color=C_WHITE)

box(s, 6.75, 6.33, 6.0, 0.65, fill=RGBColor(0x1a,0x35,0x50), radius=True, line=C_ACCENT)
txb(s, "DagsHub  —  Compare runs · Roll back · Model registry",
    6.9, 6.4, 5.7, 0.5, size=12, bold=True, color=C_ACCENT)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — Deployment
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
bar(s, "Deployment  —  Docker · Render · Streamlit Cloud · DagsHub")

# 3 service cards
svcs = [
    (C_ACCENT,  "MLflow",   "Port 5000", "Experiment Tracking\n./mlruns volume"),
    (C_GREEN,   "Backend",  "Port 8000", "FastAPI  (Render.com)\nModels bundled in image"),
    (C_ORANGE,  "Frontend", "Port 8501", "Streamlit  (Streamlit Cloud)\nThai UI · depends_on backend"),
]
for i,(col,name,port,desc) in enumerate(svcs):
    px = 0.35 + i*4.35
    box(s, px, 1.35, 4.15, 2.75, fill=C_CARD, radius=True, line=col, lw=Pt(1.5))
    txb(s, name, px+0.15, 1.45, 3.85, 0.55, size=22, bold=True, color=col)
    txb(s, port, px+0.15, 1.98, 3.85, 0.38, size=15, color=C_YELLOW)
    txb(s, desc, px+0.15, 2.4,  3.85, 1.25, size=14, color=C_LIGHT)
    if i<2:
        txb(s, "→", px+4.22, 2.48, 0.55, 0.55, size=24, color=C_LIGHT, align=PP_ALIGN.CENTER)

div(s, 4.3)

# Cloud steps
txb(s, "Cloud Deployment", 0.4, 4.45, 6.0, 0.42, size=18, bold=True, color=C_YELLOW)
steps2 = [
    (C_ACCENT,  "1", "Push to GitHub  →  Render.com auto-deploy FastAPI"),
    (C_GREEN,   "2", "Streamlit Cloud connect GitHub  →  auto-deploy frontend"),
    (C_YELLOW,  "3", "Set MLFLOW_TRACKING_URI  →  DagsHub remote tracking"),
    (C_ORANGE,  "4", "Set API_URL in Streamlit Cloud  →  point to Render backend"),
]
for i,(col,n,step) in enumerate(steps2):
    py = 5.0 + i*0.45
    box(s, 0.4, py, 0.42, 0.38, fill=col, radius=True)
    txb(s, n,    0.4, py, 0.42, 0.38, size=13, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txb(s, step, 0.95, py, 7.5, 0.38, size=13, color=C_WHITE)

# Env vars (right)
box(s, 8.7, 4.3, 4.5, 2.85, fill=C_CARD, radius=True)
txb(s, "Environment Variables", 8.85, 4.38, 4.2, 0.42, size=14, bold=True, color=C_ACCENT)
envs = [("MLFLOW_TRACKING_URI","DagsHub URL"),
        ("MLFLOW_TRACKING_USERNAME","DagsHub user"),
        ("MLFLOW_TRACKING_PASSWORD","DagsHub token"),
        ("API_URL","Render.com URL")]
for i,(k,v) in enumerate(envs):
    py = 4.88 + i*0.52
    txb(s, k, 8.85, py, 3.0, 0.38, size=11, bold=True, color=C_ACCENT)
    txb(s, v, 8.85, py+0.22, 3.0, 0.28, size=11, color=C_LIGHT)

box(s, 8.85, 7.08, 4.2, 0.35, fill=RGBColor(0x10,0x22,0x34), radius=True)
txb(s, "docker-compose up --build", 8.85, 7.08, 4.2, 0.35, size=12, color=C_YELLOW, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Results
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
bar(s, "Results  —  ผลการทดสอบโมเดล",
    "Validation set · 80/20 chronological split")

kpis = [
    (C_ACCENT,  "MAE ~0.10",    "Normalized MAE\nvalidation set"),
    (C_GREEN,   "< 16%",        "MAE สูงสุดของ\nทุก item"),
    (C_YELLOW,  "vs ~0.17",     "Baseline (predict\nmean every day)"),
    (C_ORANGE,  "MLP > Linear", "Non-linear patterns\nถูกจับได้"),
]
for i,(col,val,desc) in enumerate(kpis):
    px = 0.35 + i*3.24
    box(s, px, 1.35, 3.08, 1.8, fill=C_CARD, radius=True, line=col, lw=Pt(1.5))
    txb(s, val,  px+0.1, 1.45, 2.88, 0.72, size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, desc, px+0.1, 2.18, 2.88, 0.72, size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

div(s, 3.35)

# Per-item table (compact)
box(s, 0.35, 3.5, 6.2, 3.65, fill=RGBColor(0x10,0x22,0x34), radius=True, line=C_ACCENT, lw=Pt(1))
txb(s, "Per-Item MAE", 0.5, 3.57, 5.9, 0.42, size=15, bold=True, color=C_ACCENT)
hdrs = ["Item","Avg/day","MAE%",""]
hx   = [0.52, 2.15, 3.5, 4.8]
for j,h in enumerate(hdrs):
    txb(s, h, hx[j], 4.05, 1.55, 0.32, size=12, bold=True, color=C_YELLOW)

evals = [("Coffee","34.6","~8%",C_GREEN),("Tea","9.1","~12%",C_ACCENT),
         ("Bread","20.9","~9%",C_GREEN),("Sandwich","5.0","~14%",C_ACCENT),
         ("Cake","7.0","~13%",C_ACCENT),("Pastry","5.7","~14%",C_ACCENT),
         ("Medialuna","4.7","~15%",C_ACCENT),("Cookies","3.8","~16%",C_ORANGE),
         ("Hot Choc.","4.0","~15%",C_ACCENT),("Brownie","4.5","~16%",C_ORANGE)]
for i,(it,avg,pct,col) in enumerate(evals):
    py = 4.42 + i*0.27
    for j,v in enumerate([it,avg,pct]):
        txb(s, v, hx[j], py, 1.55, 0.25, size=11, color=C_WHITE)
    rating = "Good" if col==C_GREEN else ("Marg." if col==C_ORANGE else "OK")
    txb(s, rating, hx[3], py, 1.3, 0.25, size=11, bold=True, color=col)

# Prediction examples (right)
box(s, 6.8, 3.5, 6.15, 3.65, fill=RGBColor(0x10,0x22,0x34), radius=True, line=C_YELLOW, lw=Pt(1))
txb(s, "ตัวอย่าง — วันจันทร์ มีนาคม 2026", 6.95, 3.57, 5.85, 0.42, size=14, bold=True, color=C_YELLOW)
ehdrs = ["Item","Predict","Recommend","Waste kg"]
ehx   = [6.95, 8.55, 10.05, 11.5]
for j,h in enumerate(ehdrs):
    txb(s, h, ehx[j], 4.05, 1.45, 0.32, size=12, bold=True, color=C_ACCENT)
examples = [("Coffee","33.8","35.5","0.60"),("Tea","8.7","9.1","0.14"),
            ("Bread","20.1","21.1","0.35"),("Sandwich","4.8","5.0","0.07"),
            ("Cake","6.7","7.0","0.10"),("Pastry","5.5","5.8","0.10"),
            ("Medialuna","4.5","4.7","0.07"),("Cookies","3.7","3.9","0.07"),
            ("Hot Choc.","3.9","4.1","0.07"),("Brownie","4.3","4.5","0.07")]
for i,(it,p,r,w) in enumerate(examples):
    py = 4.42 + i*0.27
    for j,v in enumerate([it,p,r,w]):
        c = C_GREEN if j==3 else C_WHITE
        txb(s, v, ehx[j], py, 1.45, 0.25, size=11, color=c)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Impact & SDG
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)
bar(s, "Impact & Sustainability  —  SDG 2 · SDG 12")

# SDG cards
for i,(sdg,col,title,desc) in enumerate([
    ("SDG 2",  C_GREEN,   "Zero Hunger",
     "+5% Safety Stock ให้อาหารพร้อมเสมอ\nของเหลือสามารถบริจาคได้"),
    ("SDG 12", C_YELLOW, "Responsible Consumption",
     "เปลี่ยนจาก gut feeling → Data-Driven\nแสดงขยะ (กก.) สร้าง awareness"),
]):
    px = 0.35 + i*6.5
    box(s, px, 1.35, 6.2, 2.15, fill=C_CARD, radius=True, line=col, lw=Pt(2))
    box(s, px+0.15, 1.5, 1.4, 0.68, fill=col, radius=True)
    txb(s, sdg, px+0.15, 1.5, 1.4, 0.68, size=17, bold=True, color=C_BG, align=PP_ALIGN.CENTER)
    txb(s, title, px+1.65, 1.55, 4.4, 0.52, size=17, bold=True, color=col)
    txb(s, desc, px+0.2, 2.2, 5.85, 1.05, size=14, color=C_LIGHT)

div(s, 3.7)
txb(s, "ผลกระทบที่วัดได้", 0.4, 3.85, 6.0, 0.42, size=18, bold=True, color=C_ACCENT)

impacts = [
    (C_GREEN,   "4,562 กก. CO2e", "ลดต่อปี\nต่อ 1 ร้าน"),
    (C_ORANGE,  "60,000 บาท",     "ประหยัดต่อปี\nต่อ 1 ร้าน"),
    (C_ACCENT,  "< 16% MAE",      "ทุก item\nอยู่ในเกณฑ์"),
    (C_RED,     "50,000+ ร้าน",   "ตลาดเป้าหมาย\nในไทย"),
]
for i,(col,val,desc) in enumerate(impacts):
    px = 0.35 + i*3.24
    box(s, px, 4.38, 3.08, 2.25, fill=C_CARD, radius=True)
    txb(s, val,  px+0.1, 4.48, 2.88, 0.72, size=20, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, desc, px+0.1, 5.2,  2.88, 1.1,  size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

div(s, 6.75)
roadmap = [("M3","Auto-retrain"),("M6","A/B Testing"),("Y1","Multi-location"),
           ("Y2","Event features"),("Y3","AI Assistant")]
for i,(t,a) in enumerate(roadmap):
    px = 0.35+i*2.55
    box(s, px, 6.87, 2.38, 0.52, fill=C_CARD, radius=True, line=C_ACCENT, lw=Pt(1))
    txb(s, t+" · "+a, px+0.1, 6.9, 2.18, 0.46, size=12, color=C_LIGHT, align=PP_ALIGN.CENTER)


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Thank You
# ══════════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK); add_bg(s)

box(s, 9.5, -0.5, 4.5, 4.5, fill=C_CARD, radius=True)
txb(s, "🍱", 10.5, 0.5, 2.2, 1.8, size=80, align=PP_ALIGN.CENTER)

txb(s, "ขอบคุณ", 0.5, 1.5, 9.5, 1.2, size=68, bold=True)
txb(s, "Thank You", 0.5, 2.65, 9.5, 0.8, size=38, italic=True, color=C_ACCENT)

div(s, 3.65)

summary = [
    (C_ACCENT,  "20,507",       "transactions\nreal data"),
    (C_YELLOW,  "4→32→16→1",    "MLP Neural\nNetwork"),
    (C_GREEN,   "+5% Safety",   "Stock buffer\nauto-applied"),
    (C_ORANGE,  "< 16% MAE",    "All items\nin target"),
    (C_RED,     "SDG 2 + 12",   "Zero Hunger\nResp. Consumption"),
]
for i,(col,val,desc) in enumerate(summary):
    px = 0.35+i*2.55
    box(s, px, 3.88, 2.38, 2.15, fill=C_CARD, radius=True)
    txb(s, val,  px+0.1, 3.98, 2.18, 0.72, size=18, bold=True, color=col, align=PP_ALIGN.CENTER)
    txb(s, desc, px+0.1, 4.72, 2.18, 1.0,  size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)

div(s, 6.15)
txb(s, "PyTorch  ·  Prefect  ·  MLflow  ·  DagsHub  ·  FastAPI  ·  Streamlit  ·  Docker",
    0.4, 6.3, 12.5, 0.42, size=13, color=C_LIGHT, align=PP_ALIGN.CENTER)
txb(s, "Palus Kaewaram (66070131)  ·  Natthapol Aiemburanont (66070062)",
    0.4, 6.82, 12.5, 0.4, size=14, bold=True, color=C_ACCENT, align=PP_ALIGN.CENTER)


# ── Save ───────────────────────────────────────────────────────────────────────
out = "Food_Waste_Slides_v2.pptx"
prs.save(out)
print(f"Saved: {out}  ({len(prs.slides)} slides)")
